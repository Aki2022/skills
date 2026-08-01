#!/usr/bin/env python3
"""
build_template_v3.py — 旧会社テンプレ（124レイアウト・navy様式）から v3 注入用テンプレを生成する。

入力: 旧 template.pptx（1マスター・124レイアウト・Office 2007-2010テーマ・Meiryo/Calibri・#4F4F70）
出力: template_v3.pptx（10レイアウト・v3グレー体系・Noto Sans CJK JP・dk2=44546A）

変換内容（2026-07-13 設計。実測根拠は test3.pptx / tokens.json v3）:
  1. レイアウト間引き: KEEP_RENAME の10枚だけ残し改名（tone別・Preso系・content別変種は削除）
  2. テーマ: clrScheme を Office 標準系（dk2=44546A）へ、fontScheme を Noto Sans CJK JP へ
  3. 色置換: #4F4F70 → 図形塗り=44546A（positive）/ 文字色=404040（neutral）
  4. フォント: Meiryo / Arial / System Font Regular → Noto Sans CJK JP（buFontは除く）
  5. キーメッセージph(BODY idx=13): 28pt に統一（04_TitleOnly の36ptも）
  6. フッター(dt/sldNum/ftr ph): y=18.09cm・h=0.9cm・下揃え・12pt に正規化（色は既存 bg1-85%=D9D9D9）
  7. コピーライト: 「Copyright© 2026 オフィスオハナ合同会社 All Rights Reserved」へ（マスター＋全レイアウト）
  ※ 幾何（タイトル0.42/0.65・キーメッセージ0.42/1.40等）は元テンプレが既に v3 実測値と一致
    しているため変更しない。本編スライドは残す（レイアウト参照が KEEP 内なのでそのまま有効）。

使い方:
  python3 build_template_v3.py <入力template.pptx> <出力template_v3.pptx>
実行後は sanitize_pptx.py --check と PowerPoint での repair なし確認を行うこと。
"""
import re
import sys
import zipfile
import shutil
import tempfile

from pptx import Presentation
from pptx.oxml.ns import qn

A = "http://schemas.openxmlformats.org/drawingml/2006/main"

KEEP_RENAME = {
    "00_Title": "00_Title",
    "01_Contents": "01_Contents",
    "02_Section": "02_Section",
    "Handout_Single_Object_Pos": "10_Single",
    "Handout_TwoCol_Object_Pos": "11_TwoCol",
    "Handout_ThreeCol_Object_Pos": "12_ThreeCol",
    "Handout_Comparison_Object_Pos": "13_Comparison2col",
    "Handout_ThreeCol_Comparison_Object_Pos": "13_Comparison3col",
    "Handout_ObjectText_Pos": "14_ContentText",
    "04_TitleOnly": "15_TitleOnly",
    "05_Blank_Neg": "16_Blank",
}

COPYRIGHT = "Copyright© 2026 オフィスオハナ合同会社 All Rights Reserved"
FONT = "Noto Sans CJK JP"
REPLACE_FONTS = {"Meiryo", "Arial", "System Font Regular", "ＭＳ Ｐゴシック", "MS PGothic"}

# tokens.json v3
FILL_ACCENT = "44546A"   # 図形塗りの旧navy → positive slate
TEXT_NEUTRAL = "404040"  # 文字色の旧navy → neutral dark gray
FOOTER_Y = 6512560       # 18.09cm
FOOTER_CY = 323850       # 0.90cm

# Office 標準 clrScheme（test3.pptx theme 実測と同一）
CLR_SCHEME = (
    '<a:clrScheme xmlns:a="%s" name="Office">'
    '<a:dk1><a:sysClr val="windowText" lastClr="000000"/></a:dk1>'
    '<a:lt1><a:sysClr val="window" lastClr="FFFFFF"/></a:lt1>'
    '<a:dk2><a:srgbClr val="44546A"/></a:dk2>'
    '<a:lt2><a:srgbClr val="E7E6E6"/></a:lt2>'
    '<a:accent1><a:srgbClr val="4472C4"/></a:accent1>'
    '<a:accent2><a:srgbClr val="ED7D31"/></a:accent2>'
    '<a:accent3><a:srgbClr val="A5A5A5"/></a:accent3>'
    '<a:accent4><a:srgbClr val="FFC000"/></a:accent4>'
    '<a:accent5><a:srgbClr val="5B9BD5"/></a:accent5>'
    '<a:accent6><a:srgbClr val="70AD47"/></a:accent6>'
    '<a:hlink><a:srgbClr val="0563C1"/></a:hlink>'
    '<a:folHlink><a:srgbClr val="954F72"/></a:folHlink>'
    "</a:clrScheme>"
) % A


def fix_fonts(root):
    n = 0
    for tag in ("latin", "ea", "cs"):
        for el in root.iter(f"{{{A}}}{tag}"):
            tf = el.get("typeface", "")
            if tf in REPLACE_FONTS or (tag == "ea" and tf == ""):
                # buFont は対象外（親タグで判定）
                el.set("typeface", FONT)
                for junk in ("panose", "pitchFamily", "charset"):
                    if junk in el.attrib:
                        del el.attrib[junk]
                n += 1
    return n


def fix_colors(root):
    """srgbClr 4F4F70 を文脈で置換: 図形塗り(spPrの子)=44546A / 文字(rPr系)=404040"""
    TEXT_CTX = {f"{{{A}}}defRPr", f"{{{A}}}rPr", f"{{{A}}}endParaRPr"}
    parent_map = {c: p for p in root.iter() for c in p}
    n = 0
    for clr in list(root.iter(f"{{{A}}}srgbClr")):
        if clr.get("val", "").upper() != "4F4F70":
            continue
        anc, is_text = clr, False
        while anc in parent_map:
            anc = parent_map[anc]
            if anc.tag in TEXT_CTX:
                is_text = True
                break
        clr.set("val", TEXT_NEUTRAL if is_text else FILL_ACCENT)
        n += 1
    return n


def fix_copyright(shapes):
    n = 0
    for sh in shapes:
        if not sh.has_text_frame:
            continue
        tf = sh.text_frame
        if not tf.text.startswith("Copyright©"):
            continue
        para = tf.paragraphs[0]
        runs = para.runs
        if runs:
            runs[0].text = COPYRIGHT
            for r in runs[1:]:
                r._r.getparent().remove(r._r)
            n += 1
    return n


def fix_footer_phs(shapes):
    """dt/sldNum/ftr プレースホルダを v3 位置・サイズ・12pt・下揃えに正規化"""
    n = 0
    for sh in shapes:
        if not sh.is_placeholder:
            continue
        ph_type = sh.element.find(f".//{{{'http://schemas.openxmlformats.org/presentationml/2006/main'}}}ph")
        t = ph_type.get("type") if ph_type is not None else None
        if t not in ("dt", "sldNum", "ftr"):
            continue
        el = sh.element
        xfrm = el.find(f".//{{{A}}}xfrm")
        if xfrm is not None:
            off = xfrm.find(f"{{{A}}}off")
            ext = xfrm.find(f"{{{A}}}ext")
            if off is not None and ext is not None:
                off.set("y", str(FOOTER_Y))
                ext.set("cy", str(FOOTER_CY))
        body = el.find(f".//{{{A}}}bodyPr")
        if body is not None:
            body.set("anchor", "b")
        for pr in el.iter(f"{{{A}}}defRPr"):
            pr.set("sz", "1200")
        for pr in el.iter(f"{{{A}}}endParaRPr"):
            if pr.get("sz"):
                pr.set("sz", "1200")
        n += 1
    return n


KEYMSG_LAYOUTS = {
    "10_Single",
    "11_TwoCol",
    "12_ThreeCol",
    "13_Comparison2col",
    "13_Comparison3col",
    "14_ContentText",
    "15_TitleOnly",
}


def fix_keymsg(layout, new_name):
    """キーメッセージph(BODY idx=13)を28ptに統一。01_Contents の目次リストも同じ idx=13 のため、
    キーメッセージ帯を持つレイアウト（KEYMSG_LAYOUTS）に限定し、さらに位置（帯は y≤2cm）でも絞る。"""
    if new_name not in KEYMSG_LAYOUTS:
        return 0
    P = "http://schemas.openxmlformats.org/presentationml/2006/main"
    n = 0
    for sh in layout.shapes:
        if not sh.is_placeholder:
            continue
        ph = sh.element.find(f".//{{{P}}}ph")
        if ph is None or ph.get("idx") != "13" or ph.get("type", "body") not in ("body", None):
            continue
        if sh.top is None or sh.top > 720000:  # 2cm より下は本文リスト
            continue
        for pr in sh.element.iter(f"{{{A}}}defRPr"):
            pr.set("sz", "2800")
            n += 1
    return n


def main():
    if len(sys.argv) != 3:
        print(__doc__)
        sys.exit(2)
    src, dst = sys.argv[1], sys.argv[2]
    prs = Presentation(src)
    master = prs.slide_masters[0]

    # 1. レイアウト間引き
    sldLayoutIdLst = master.element.find(qn("p:sldLayoutIdLst"))
    removed = 0
    for layout in list(master.slide_layouts):
        if layout.name in KEEP_RENAME:
            continue
        for slid in list(sldLayoutIdLst):
            r_id = slid.get(qn("r:id"))
            if master.part.related_part(r_id) is layout.part:
                sldLayoutIdLst.remove(slid)
                master.part.drop_rel(r_id)
                removed += 1
                break
    print(f"layouts removed: {removed}, kept: {len(list(master.slide_layouts))}")

    # 2. マスターの修正
    stats = dict(fonts=0, colors=0, cr=0, footer=0, keymsg=0)
    stats["fonts"] += fix_fonts(master.element)
    stats["colors"] += fix_colors(master.element)
    stats["cr"] += fix_copyright(master.shapes)
    stats["footer"] += fix_footer_phs(master.shapes)

    # 3. 残レイアウトの修正＋改名
    for layout in master.slide_layouts:
        new = KEEP_RENAME[layout.name]
        layout.element.find(qn("p:cSld")).set("name", new)
        stats["fonts"] += fix_fonts(layout.element)
        stats["colors"] += fix_colors(layout.element)
        stats["cr"] += fix_copyright(layout.shapes)
        stats["footer"] += fix_footer_phs(layout.shapes)
        stats["keymsg"] += fix_keymsg(layout, new)
    print("fixes:", stats)

    prs.save(dst)

    # 4. テーマ差し替え（zipレベル: clrScheme + fontScheme）
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    tmp.close()
    with zipfile.ZipFile(dst) as zin, zipfile.ZipFile(tmp.name, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == "ppt/theme/theme1.xml":
                xml = data.decode("utf-8")
                xml = re.sub(r"<a:clrScheme .*?</a:clrScheme>", CLR_SCHEME.replace(f' xmlns:a="{A}"', ""), xml, flags=re.S)
                xml = re.sub(
                    r'(<a:(?:major|minor)Font>)<a:latin typeface="[^"]*"[^/]*/><a:ea typeface="[^"]*"[^/]*/>',
                    rf'\1<a:latin typeface="{FONT}"/><a:ea typeface="{FONT}"/>',
                    xml,
                )
                data = xml.encode("utf-8")
            zout.writestr(item, data)
    shutil.move(tmp.name, dst)
    print(f"done: {dst}")


if __name__ == "__main__":
    main()
